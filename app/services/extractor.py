"""Text extraction service for various document formats."""

import asyncio
import logging
from functools import partial
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


class TextExtractor:
    """Extract text from various document formats."""

    SUPPORTED_TEXT = {".txt", ".md"}
    SUPPORTED_PDF = {".pdf"}
    SUPPORTED_PPTX = {".pptx"}
    SUPPORTED_AUDIO = {".mp3", ".wav", ".m4a"}

    def __init__(self, whisper_model: str = "large") -> None:
        self.whisper_model = whisper_model
        self._whisper: Any = None

    async def extract(self, file_path: Path) -> str:
        """Extract text from a file based on its extension."""
        suffix = file_path.suffix.lower()

        if suffix in self.SUPPORTED_TEXT:
            return await self._extract_text(file_path)
        elif suffix in self.SUPPORTED_PDF:
            return await self._extract_pdf(file_path)
        elif suffix in self.SUPPORTED_PPTX:
            return await self._extract_pptx(file_path)
        elif suffix in self.SUPPORTED_AUDIO:
            return await self._extract_audio(file_path)
        else:
            raise ValueError(f"Unsupported file format: {suffix}")

    async def _extract_text(self, file_path: Path) -> str:
        """Extract text from plain text or markdown files."""
        logger.info(f"Extracting text from {file_path}")
        # Run file I/O in thread executor to avoid blocking
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, partial(file_path.read_text, encoding="utf-8"))

    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files using pdfplumber."""
        logger.info(f"Extracting text from PDF {file_path}")
        # Run blocking PDF extraction in thread executor
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._extract_pdf_sync, file_path)

    def _extract_pdf_sync(self, file_path: Path) -> str:
        """Synchronous PDF extraction for use in executor."""
        import pdfplumber

        text_parts: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n\n".join(text_parts)

    async def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint files including speaker notes."""
        logger.info(f"Extracting text from PPTX {file_path}")
        # Run blocking PPTX extraction in thread executor
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._extract_pptx_sync, file_path)

    def _extract_pptx_sync(self, file_path: Path) -> str:
        """Synchronous PPTX extraction for use in executor."""
        from pptx import Presentation

        text_parts: list[str] = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame.text:
                    slide_text.append(f"[Notes: {notes_frame.text}]")

            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

        return "\n\n".join(text_parts)

    async def _extract_audio(self, file_path: Path) -> str:
        """Transcribe audio files using Whisper."""
        logger.info(f"Transcribing audio {file_path} with Whisper {self.whisper_model}")
        # Run CPU-intensive Whisper transcription in thread executor
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._transcribe_audio_sync, file_path)

    def _transcribe_audio_sync(self, file_path: Path) -> str:
        """Synchronous audio transcription for use in executor."""
        import whisper

        # Load model lazily
        if self._whisper is None:
            logger.info(f"Loading Whisper model: {self.whisper_model}")
            self._whisper = whisper.load_model(self.whisper_model)

        result: dict[str, Any] = self._whisper.transcribe(str(file_path), language="de")
        text: str = result["text"]
        return text

    @classmethod
    def supported_extensions(cls) -> set[str]:
        """Get all supported file extensions."""
        return cls.SUPPORTED_TEXT | cls.SUPPORTED_PDF | cls.SUPPORTED_PPTX | cls.SUPPORTED_AUDIO
